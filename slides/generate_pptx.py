#!/usr/bin/env python3
"""Generate Good Vibes ECCB 2026 slides as a .pptx file with ZSL branding.

Upload the output to Google Drive and open as Google Slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Brand colours ──
BURNHAM = RGBColor(0x00, 0x21, 0x10)
FUN_GREEN = RGBColor(0x00, 0x89, 0x45)
MID_GREEN = RGBColor(0x00, 0x4D, 0x25)
SAPLING = RGBColor(0xDE, 0xC6, 0xA2)
ELECTRIC_LIME = RGBColor(0xC7, 0xF7, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ERROR_RED = RGBColor(0xFF, 0x6B, 0x6B)
LIGHT_WHITE = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use the blank layout
blank_layout = prs.slide_layouts[6]

LOGO_PATH = os.path.join(os.path.dirname(__file__), "assets", "zsl-ioz.png")
HAS_LOGO = os.path.exists(LOGO_PATH)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide):
    if not HAS_LOGO:
        return
    slide.shapes.add_picture(
        LOGO_PATH,
        SLIDE_WIDTH - Inches(1.8),  # right-aligned
        Inches(0.25),
        width=Inches(1.5),
    )


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Arial"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False,
             alignment=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Arial"
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


# ════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
            "Good Vibes", font_size=60, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(3.0), Inches(11), Inches(1),
            "Responsible use of AI for large-scale,\nrapid Earth Observation workflows",
            font_size=24, color=SAPLING, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
            "AB  ·  HFT  ·  HH  ·  JW",
            font_size=20, color=LIGHT_WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(4.8), Inches(11), Inches(0.5),
            "Institute of Zoology, Zoological Society of London",
            font_size=16, color=LIGHT_WHITE, alignment=PP_ALIGN.CENTER)

# ECCB badge
badge = s.shapes.add_textbox(Inches(5.5), Inches(5.6), Inches(2.3), Inches(0.5))
bf = badge.text_frame
bp = bf.paragraphs[0]
bp.text = "ECCB 2026"
bp.font.size = Pt(18)
bp.font.bold = True
bp.font.color.rgb = WHITE
bp.font.name = "Arial"
bp.alignment = PP_ALIGN.CENTER
badge.fill.solid()
badge.fill.fore_color.rgb = FUN_GREEN

add_notes(s, "Welcome everyone to 'Good Vibes'. This workshop is about using AI responsibly — specifically for Earth Observation coding workflows. We'll cover what AI is, why it matters for conservation science, and six principles for working with it safely. Then we'll get hands-on.")


# ════════════════════════════════════════
# SLIDE 2: SECTION HEADER — What is AI
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, MID_GREEN)
add_logo(s)

add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
            "🤖 What is AI?", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
            "…and why should you care?",
            font_size=24, color=SAPLING, alignment=PP_ALIGN.CENTER)

add_notes(s, "Let's start with the basics. We're going to take about 5 minutes to frame what AI is, how we should think about it, and what it can do for us as researchers.")


# ════════════════════════════════════════
# SLIDE 3: The ChatGPT Moment
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "The ChatGPT Moment", font_size=40, color=WHITE, bold=True)
add_textbox(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "An explosion of capability since Nov 2022",
            font_size=22, color=ELECTRIC_LIME)

timeline = [
    ("Nov 2022", "ChatGPT launches → 100M users in 2 months"),
    ("2023", "GPT-4, Bard, Claude — multi-modal capabilities"),
    ("2024", "Reasoning models, agentic AI emerges"),
    ("2025", "AI coding agents, science-specific models"),
    ("2026", "📍 You are here — AI in conservation workflows"),
]

y = 2.2
for year, desc in timeline:
    tf = add_textbox(s, Inches(1.5), Inches(y), Inches(2), Inches(0.4),
                     year, font_size=18, color=ELECTRIC_LIME, bold=True)
    add_textbox(s, Inches(3.8), Inches(y), Inches(8), Inches(0.4),
                desc, font_size=17, color=WHITE)
    y += 0.7

add_notes(s, "ChatGPT launched in November 2022 and reached 100 million users in just two months — the fastest adoption of any technology in history. Since then, the pace has been extraordinary. We've gone from novelty chatbots to multi-modal reasoning systems that can write code, analyse data, and even operate autonomously.")


# ════════════════════════════════════════
# SLIDE 4: How We Think About AI
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "How We Think About AI", font_size=40, color=WHITE, bold=True)
add_textbox(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "The metaphor spectrum", font_size=22, color=ELECTRIC_LIME)

metaphors = ["🧑‍🤝‍🧑  Colleague", "🤝  Peer", "🔧  Tool", "🔍  Search Engine"]
x = 1.0
for m in metaphors:
    box = s.shapes.add_textbox(Inches(x), Inches(2.5), Inches(2.6), Inches(1.0))
    bf = box.text_frame
    bp = bf.paragraphs[0]
    bp.text = m
    bp.font.size = Pt(20)
    bp.font.color.rgb = WHITE
    bp.font.bold = True
    bp.font.name = "Arial"
    bp.alignment = PP_ALIGN.CENTER
    box.fill.solid()
    box.fill.fore_color.rgb = FUN_GREEN
    x += 2.9

tf = add_textbox(s, Inches(1.5), Inches(4.2), Inches(10), Inches(1.2),
                 "⚠️  AI is none of these perfectly.", font_size=20,
                 color=SAPLING, alignment=PP_ALIGN.CENTER, bold=True)
add_para(tf, "It reasons — but not like a human. It retrieves — but not like a database.",
         font_size=18, color=LIGHT_WHITE, alignment=PP_ALIGN.CENTER)
add_para(tf, "How you think about it shapes how you use it.",
         font_size=18, color=LIGHT_WHITE, alignment=PP_ALIGN.CENTER)

add_notes(s, "People tend to reach for familiar metaphors when thinking about AI. Some treat it like a colleague — expecting judgement and accountability. Others treat it like a search engine — expecting it to just look things up. The truth is it's none of these perfectly.")


# ════════════════════════════════════════
# SLIDE 5: Use 1 — Reviewing & Conceptualisation
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "📚  Use 1: Reviewing & Conceptualisation",
            font_size=36, color=WHITE, bold=True)

uses = [
    ("🔬", "Literature Review", "Synthesise papers, identify gaps, find connections across disciplines"),
    ("💡", "Brainstorming", "Generate hypotheses, explore alternative framings, challenge assumptions"),
    ("🧪", "Hypothesis Generation", "Identify testable predictions from complex datasets and theory"),
]

y = 2.0
for emoji, title, desc in uses:
    add_textbox(s, Inches(1.2), Inches(y), Inches(0.8), Inches(0.6),
                emoji, font_size=28, alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.2), Inches(y), Inches(4), Inches(0.5),
                title, font_size=22, color=ELECTRIC_LIME, bold=True)
    add_textbox(s, Inches(2.2), Inches(y + 0.45), Inches(9), Inches(0.5),
                desc, font_size=17, color=LIGHT_WHITE)
    y += 1.3

add_notes(s, "The first big use case: conceptual work. AI is genuinely useful for literature review — not as a replacement for reading papers, but for synthesising across large bodies of work, spotting gaps, and helping you think across disciplines.")


# ════════════════════════════════════════
# SLIDE 6: Use 2 — Maths & Reasoning
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "🧮  Use 2: Maths & Reasoning",
            font_size=36, color=WHITE, bold=True)

uses = [
    ("📊", "Statistical Analysis", "Model selection, power calculations, interpreting complex outputs"),
    ("📈", "Data Interpretation", "Explaining results, identifying patterns, suggesting visualisations"),
    ("🔢", "Complex Calculations", "Unit conversions, spatial statistics, remote sensing indices"),
]

y = 2.0
for emoji, title, desc in uses:
    add_textbox(s, Inches(1.2), Inches(y), Inches(0.8), Inches(0.6),
                emoji, font_size=28, alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.2), Inches(y), Inches(4), Inches(0.5),
                title, font_size=22, color=ELECTRIC_LIME, bold=True)
    add_textbox(s, Inches(2.2), Inches(y + 0.45), Inches(9), Inches(0.5),
                desc, font_size=17, color=LIGHT_WHITE)
    y += 1.3

add_notes(s, "The second use case: quantitative reasoning. Modern LLMs are surprisingly capable at mathematics — model selection advice, interpreting statistical outputs, even working through complex derivations.")


# ════════════════════════════════════════
# SLIDE 7: Use 3 — Coding
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            "💻  Use 3: Coding",
            font_size=36, color=WHITE, bold=True)
add_textbox(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "This is what we'll focus on today",
            font_size=22, color=ELECTRIC_LIME)

uses = [
    ("✍️", "Writing Code", "Generate scripts from natural language descriptions of your analysis"),
    ("🐛", "Debugging", "Identify bugs, explain errors, suggest fixes — often faster than Stack Overflow"),
    ("♻️", "Refactoring", "Improve code structure, add documentation, optimise performance"),
]

y = 2.2
for emoji, title, desc in uses:
    add_textbox(s, Inches(1.2), Inches(y), Inches(0.8), Inches(0.6),
                emoji, font_size=28, alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.2), Inches(y), Inches(4), Inches(0.5),
                title, font_size=22, color=ELECTRIC_LIME, bold=True)
    add_textbox(s, Inches(2.2), Inches(y + 0.45), Inches(9), Inches(0.5),
                desc, font_size=17, color=LIGHT_WHITE)
    y += 1.3

add_notes(s, "And here we get to today's focus: coding. AI can write code from scratch based on your descriptions, debug existing code by explaining errors and suggesting fixes, and refactor messy scripts into clean, documented, well-structured code.")


# ════════════════════════════════════════
# SLIDE 8: What is Vibe Coding?
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
            '🎵  What is "Vibe Coding"?',
            font_size=36, color=WHITE, bold=True)

# Quote
quote_box = s.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(1.5))
qf = quote_box.text_frame
qf.word_wrap = True
qp = qf.paragraphs[0]
qp.text = '"There\'s a new kind of coding I call vibe coding, where you fully give in to the vibes, embrace exponentials, and forget that the code even exists."'
qp.font.size = Pt(18)
qp.font.color.rgb = SAPLING
qp.font.italic = True
qp.font.name = "Arial"
add_para(qf, "— Andrej Karpathy, Feb 2025", font_size=14, color=LIGHT_WHITE)

# Spectrum
labels = [
    ("Fully manual", FUN_GREEN),
    ("AI-assisted", RGBColor(0x4D, 0xB8, 0x78)),
    ("AI-led", ELECTRIC_LIME),
    ("Fully AI-generated", ERROR_RED),
]
x = 1.5
for label, clr in labels:
    box = s.shapes.add_textbox(Inches(x), Inches(4.2), Inches(2.3), Inches(0.8))
    bf = box.text_frame
    bp = bf.paragraphs[0]
    bp.text = label
    bp.font.size = Pt(16)
    bp.font.color.rgb = WHITE
    bp.font.bold = True
    bp.font.name = "Arial"
    bp.alignment = PP_ALIGN.CENTER
    box.fill.solid()
    box.fill.fore_color.rgb = clr
    x += 2.6

add_textbox(s, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
            "For scientific code, we need to stay in the responsible middle of this spectrum.",
            font_size=18, color=SAPLING, alignment=PP_ALIGN.CENTER)

add_notes(s, "Andrej Karpathy coined the term 'vibe coding' in early 2025. It describes the experience of describing what you want in natural language and letting AI write all the code. For scientific code — where policy decisions depend on correct analysis — we need to be somewhere in the responsible middle.")


# ════════════════════════════════════════
# SLIDE 9: Transition — "responsibly?"
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(1), Inches(2.0), Inches(11), Inches(0.8),
            "But how do we do this", font_size=28, color=SAPLING,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(2.8), Inches(11), Inches(1.2),
            "responsibly?", font_size=60, color=ELECTRIC_LIME, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
            "↓", font_size=48, color=ELECTRIC_LIME,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
            "Six Principles for Working with AI",
            font_size=20, color=LIGHT_WHITE, alignment=PP_ALIGN.CENTER)

add_notes(s, "So we've seen what AI can do. But the critical question is: how do we use these capabilities responsibly? That brings us to our six principles.")


# ════════════════════════════════════════
# SLIDE 10: SECTION HEADER — Six Principles
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, MID_GREEN)
add_logo(s)

add_textbox(s, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
            "📋  Six Principles", font_size=48, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
            "for Responsible AI-Assisted Coding",
            font_size=24, color=SAPLING, alignment=PP_ALIGN.CENTER)

add_notes(s, "This is the core of today's framing talk. Six principles that we think every researcher should follow when using AI to write code.")


# ────────────────────────────────────────
# Helper for principle slides
# ────────────────────────────────────────
def make_principle_slide(num, emoji, title, tagline, bullets, note, tip=None):
    s = prs.slides.add_slide(blank_layout)
    set_bg(s, BURNHAM)
    add_logo(s)

    add_textbox(s, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
                f"Principle {num}", font_size=20, color=LIGHT_WHITE)

    add_textbox(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
                f"{emoji}  {title}", font_size=36, color=WHITE, bold=True)

    # Tagline card
    card = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.5), Inches(0.7))
    cf = card.text_frame
    cp = cf.paragraphs[0]
    cp.text = tagline
    cp.font.size = Pt(20)
    cp.font.color.rgb = SAPLING
    cp.font.name = "Arial"
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x00, 0x35, 0x1A)

    y = 3.0
    for b in bullets:
        tf = add_textbox(s, Inches(1.2), Inches(y), Inches(10.5), Inches(0.55),
                         "", font_size=17, color=WHITE)
        p = tf.paragraphs[0]
        # Split on first " — " to bold the first part
        if " — " in b:
            parts = b.split(" — ", 1)
            run1 = p.add_run()
            run1.text = parts[0] + " — "
            run1.font.bold = True
            run1.font.size = Pt(17)
            run1.font.color.rgb = ELECTRIC_LIME
            run1.font.name = "Arial"
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(17)
            run2.font.color.rgb = WHITE
            run2.font.name = "Arial"
        else:
            p.text = b
            p.font.size = Pt(17)
            p.font.color.rgb = WHITE
            p.font.name = "Arial"
        y += 0.65

    if tip:
        add_textbox(s, Inches(1.5), Inches(y + 0.2), Inches(10), Inches(0.6),
                    tip, font_size=16, color=ELECTRIC_LIME,
                    alignment=PP_ALIGN.CENTER)

    add_notes(s, note)
    return s


# SLIDES 11-16: Principles 1-6
make_principle_slide(1, "🧠", "Treat It as an Intelligence",
    "It reasons. Give it context, domain knowledge, and specifics.",
    [
        'Provide domain context — "I\'m analysing Sentinel-2 imagery for land cover classification in tropical forests"',
        'Be specific — "Use the SCL band to mask clouds, not QA60" rather than "remove clouds"',
        "Share your constraints — memory limits, time series length, study area extent",
        "Don't treat it like a search bar — it can reason through multi-step problems if you frame them well",
    ],
    "Principle 1: treat AI as an intelligence, not a search engine. It can reason — but only if you give it the right context. Tell it what satellite you're using, what region, what time period, what your scientific question is.")

make_principle_slide(2, "🔍", "Distrust by Default",
    "Assume output is plausible but wrong until verified.",
    [
        "It will invent — function names, band names, and API calls that don't exist",
        "It will produce code that looks correct — but has subtle logical errors",
        "It will present wrong answers — with complete confidence",
    ],
    "Principle 2: distrust by default. This is perhaps the most important principle. AI is confidently incorrect. It will generate code that looks perfectly reasonable, uses plausible-sounding function names, and produces output — but the output might be completely wrong.",
    tip="🎭 Think of it as a very keen but unreliable junior colleague")

make_principle_slide(3, "✅", "Verify with Tests",
    "If you can't test it, you can't trust it.",
    [
        "Write unit tests — test individual functions with known inputs and expected outputs",
        "Check edge cases — empty collections, null geometries, single-pixel regions",
        "Validate against known answers — use published datasets, manual calculations, or sanity checks",
        "Ask AI to write tests too — then verify the tests themselves!",
    ],
    "Principle 3: verify with tests. Unit testing isn't just good software practice — when working with AI, it's essential.",
    tip="💡 Generate code and tests separately, then check they agree")

make_principle_slide(4, "📖", "Prioritise Reporting & Interpretation",
    "Don't just run code — understand what it does.",
    [
        "Can you explain every line? — If not, you don't understand your analysis",
        "Document assumptions — why this reducer? Why this scale? Why this time period?",
        "Interpret results critically — does the output make ecological sense?",
        "Write for your future self — add comments that explain why, not just what",
    ],
    "Principle 4: prioritise reporting and interpretation. If you can't explain every line of that code, you don't actually understand your analysis.")

make_principle_slide(5, "💥", "Engineer Adversity",
    "Deliberately try to break your AI-generated code.",
    [
        "Feed it bad inputs — wrong date formats, corrupted imagery, mismatched CRS",
        "Test edge cases — polar regions, date-line crossings, very large/small areas",
        "Try empty datasets — what happens when your filter returns zero images?",
        "If it breaks, it wasn't robust — and that's information you need before publication",
    ],
    "Principle 5: engineer adversity. Don't just test the happy path — actively try to break your code.",
    tip="🎯 Think like a hostile reviewer — what would they try?")

make_principle_slide(6, "🔓", "Maximise Transparency",
    "Document which parts are AI-generated and how.",
    [
        "Log your prompts — save the conversations that generated your code",
        "Version your code — use Git; track the evolution from AI draft to final version",
        "Capture metadata — which model, which version, what date",
        'Be honest in methods — "Code was initially generated using [model] and subsequently reviewed, tested, and modified by the authors"',
    ],
    "Principle 6: maximise transparency. This is about scientific integrity. Log your prompts, use version control, record which model you used, and be honest in your methods section.")


# ════════════════════════════════════════
# SLIDE 17: Six Principles Summary
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "The Six Principles", font_size=40, color=WHITE, bold=True)
add_textbox(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
            "Your checklist for responsible AI coding",
            font_size=22, color=ELECTRIC_LIME)

principles_summary = [
    ("🧠", "1. Treat as Intelligence", "Give context, be specific"),
    ("🔍", "2. Distrust by Default", "Assume wrong until verified"),
    ("✅", "3. Verify with Tests", "Can't test it? Can't trust it"),
    ("📖", "4. Report & Interpret", "Understand every line"),
    ("💥", "5. Engineer Adversity", "Try to break it"),
    ("🔓", "6. Maximise Transparency", "Log, version, disclose"),
]

y = 2.2
for emoji, name, reminder in principles_summary:
    add_textbox(s, Inches(1.2), Inches(y), Inches(0.8), Inches(0.5),
                emoji, font_size=24, alignment=PP_ALIGN.CENTER)
    add_textbox(s, Inches(2.2), Inches(y), Inches(4.5), Inches(0.5),
                name, font_size=20, color=WHITE, bold=True)
    add_textbox(s, Inches(7.0), Inches(y), Inches(5), Inches(0.5),
                reminder, font_size=18, color=SAPLING)
    y += 0.7

add_notes(s, "Here are all six principles together as a checklist. You might want to take a photo of this slide.")


# ════════════════════════════════════════
# SLIDE 18: Real Failures
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "⚠️  Real Failures", font_size=40, color=WHITE, bold=True)
add_textbox(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
            "Common AI mistakes in Google Earth Engine",
            font_size=22, color=ELECTRIC_LIME)

failures = [
    ("❌ Wrong Band Names",
     'AI uses B4, B8 for Sentinel-2 surface reflectance — but the actual band names are SR_B4, SR_B8. Code runs without error but produces garbage.'),
    ("❌ Incorrect Scale",
     'AI defaults to scale: 30 (Landsat) for Sentinel-2 analysis. Results are technically valid but scientifically wrong — resampled at the wrong resolution.'),
    ("❌ Misapplied Reducers",
     'AI uses ee.Reducer.mean() where ee.Reducer.median() was needed for cloud-free compositing. Or applies a spatial reducer when a temporal reducer was intended.'),
]

y = 2.0
for title, desc in failures:
    # Card background
    card = s.shapes.add_textbox(Inches(1.0), Inches(y), Inches(11.3), Inches(1.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x2D, 0x0A, 0x0A)
    cf = card.text_frame
    cf.word_wrap = True
    cp = cf.paragraphs[0]
    cp.text = title
    cp.font.size = Pt(19)
    cp.font.color.rgb = ERROR_RED
    cp.font.bold = True
    cp.font.name = "Arial"
    dp = cf.add_paragraph()
    dp.text = desc
    dp.font.size = Pt(15)
    dp.font.color.rgb = LIGHT_WHITE
    dp.font.name = "Arial"
    dp.space_before = Pt(4)
    y += 1.6

add_notes(s, "Let's make this concrete with real examples from Google Earth Engine. Wrong band names, incorrect scale, and misapplied reducers are exactly the kinds of bugs that Principles 2 through 5 are designed to catch.")


# ════════════════════════════════════════
# SLIDE 19: Why This Matters
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
            "🌍  Why This Matters", font_size=40, color=WHITE, bold=True)
add_textbox(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
            "For conservation science", font_size=22, color=ELECTRIC_LIME)

# Callout
callout = s.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.0))
callout.fill.solid()
callout.fill.fore_color.rgb = RGBColor(0x00, 0x35, 0x1A)
cf = callout.text_frame
cf.word_wrap = True
cp = cf.paragraphs[0]
cp.text = "Policy decisions depend on correct analysis.\nA bug in your deforestation script isn't just a software problem — it's a conservation problem."
cp.font.size = Pt(18)
cp.font.color.rgb = SAPLING
cp.font.name = "Arial"
cp.alignment = PP_ALIGN.CENTER

stakes = [
    "🌳  Habitat loss estimates feed into protected area designations",
    "🦁  Species distribution models inform IUCN Red List assessments",
    "🔥  Fire severity maps guide post-disaster response and funding allocation",
    "💰  Carbon stock estimates drive climate finance decisions worth millions",
]

y = 3.5
for item in stakes:
    add_textbox(s, Inches(1.5), Inches(y), Inches(10), Inches(0.5),
                item, font_size=18, color=WHITE)
    y += 0.6

add_textbox(s, Inches(1), Inches(y + 0.3), Inches(11), Inches(0.5),
            "The stakes are too high for unchecked AI-generated code.",
            font_size=20, color=ELECTRIC_LIME, bold=True,
            alignment=PP_ALIGN.CENTER)

add_notes(s, "In conservation science, the stakes are real. Habitat loss estimates feed directly into protected area designations. Species models inform Red List assessments. A subtle bug in your analysis code could lead to wrong conclusions, wrong policy decisions, and ultimately harm to the species and ecosystems we're trying to protect.")


# ════════════════════════════════════════
# SLIDE 20: Transition — into practice
# ════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
set_bg(s, BURNHAM)
add_logo(s)

add_textbox(s, Inches(1), Inches(2.0), Inches(11), Inches(0.8),
            "Now let's put this", font_size=28, color=SAPLING,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(2.8), Inches(11), Inches(1.2),
            "into practice", font_size=60, color=ELECTRIC_LIME, bold=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
            "→", font_size=48, color=ELECTRIC_LIME,
            alignment=PP_ALIGN.CENTER)

add_textbox(s, Inches(1), Inches(5.3), Inches(11), Inches(0.5),
            "Hands-on practical session",
            font_size=20, color=LIGHT_WHITE, alignment=PP_ALIGN.CENTER)

add_notes(s, "That's the framing complete. You now have a mental model for what AI can do and six principles for using it responsibly. In the next section, we'll put this into practice.")


# ── Save ──
output_path = os.path.join(os.path.dirname(__file__), "..", "Good_Vibes_ECCB2026.pptx")
prs.save(output_path)
print(f"✅ Saved to {os.path.abspath(output_path)}")
print(f"   {len(prs.slides)} slides generated")
print(f"   Upload to Google Drive → Open with Google Slides")
